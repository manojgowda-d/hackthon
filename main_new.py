"""
╔══════════════════════════════════════════════════════════════════╗
║               MINI NOTEBOOKLM — Backend (main.py)               ║
║         FastAPI + RAG + ChromaDB + Gemini/OpenAI API             ║
╚══════════════════════════════════════════════════════════════════╝
"""

# ─────────────────────────────────────────────
# IMPORTS
# ─────────────────────────────────────────────

import os
import uuid
import shutil
import logging
from pathlib import Path
from typing import Optional
from datetime import datetime
from collections import deque

import uvicorn
from dotenv import load_dotenv
from fastapi import FastAPI, File, UploadFile, HTTPException, Query
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, FileResponse
from pydantic import BaseModel

# Document loaders
from pypdf import PdfReader
import docx  # python-docx
from pptx import Presentation  # python-pptx

# Text-to-speech
from gtts import gTTS

# LangChain components
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_chroma import Chroma
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_core.documents import Document

# LLM Providers (pluggable)
import google.generativeai as genai
# from openai import OpenAI  # Uncomment to use OpenAI instead

# ─────────────────────────────────────────────
# CONFIG
# ─────────────────────────────────────────────

load_dotenv()

# Logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S",
)
logger = logging.getLogger("mini-notebooklm")

# Directories
UPLOAD_DIR = Path("uploads")
CHROMA_DIR = Path("chroma_db")
VOICE_DIR  = Path("voice_overviews")
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
CHROMA_DIR.mkdir(parents=True, exist_ok=True)
VOICE_DIR.mkdir(parents=True, exist_ok=True)

# API Keys
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "")
LLM_PROVIDER = os.getenv("LLM_PROVIDER", "gemini").lower()  # "gemini" or "openai"

# Embedding model (runs locally, no API key needed)
EMBEDDING_MODEL = os.getenv("EMBEDDING_MODEL", "all-MiniLM-L6-v2")

# RAG settings
CHUNK_SIZE = int(os.getenv("CHUNK_SIZE", "800"))
CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP", "100"))
TOP_K_CHUNKS = int(os.getenv("TOP_K_CHUNKS", "5"))

# Session memory (last N messages per session)
MAX_HISTORY = int(os.getenv("MAX_HISTORY", "10"))

# Server
HOST = os.getenv("HOST", "0.0.0.0")
PORT = int(os.getenv("PORT", "5000"))

# Allowed file types
ALLOWED_EXTENSIONS = {".pdf", ".txt", ".docx",".pptx"}

# ─────────────────────────────────────────────
# APP INIT
# ─────────────────────────────────────────────

app = FastAPI(
    title="Mini NotebookLM API",
    description="A lightweight RAG-powered document chat backend.",
    version="1.0.0",
    docs_url="/docs",
    redoc_url="/redoc",
)

# CORS — allow all origins for local dev; tighten in production
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ─────────────────────────────────────────────
# MODELS (Pydantic request/response schemas)
# ─────────────────────────────────────────────

class AskRequest(BaseModel):
    question: str
    session_id: Optional[str] = "default"
    include_sources: Optional[bool] = True

class AskResponse(BaseModel):
    answer: str
    sources: Optional[list[dict]] = []
    session_id: str

class SummarizeRequest(BaseModel):
    file_name: Optional[str] = None  # Summarize specific file or all docs

class SummarizeResponse(BaseModel):
    summary: str
    file_name: Optional[str] = None

class DeleteResponse(BaseModel):
    message: str
    file_name: str

class HealthResponse(BaseModel):
    status: str
    provider: str
    uploaded_files: int
    timestamp: str

class VoiceOverviewRequest(BaseModel):
    file_name: Optional[str] = None   # None → summarise all docs
    language: Optional[str] = "en"    # gTTS language code

class VoiceOverviewResponse(BaseModel):
    audio_file: str
    summary_text: str
    file_name: Optional[str] = None

class MindMapRequest(BaseModel):
    file_name: Optional[str] = None   # None → use all indexed docs
    output_format: Optional[str] = "mermaid"  # "mermaid" or "json"

class MindMapResponse(BaseModel):
    mind_map: str          # Mermaid diagram string OR JSON string
    output_format: str
    file_name: Optional[str] = None

# ─────────────────────────────────────────────
# UTILITY: Session Memory
# ─────────────────────────────────────────────

# In-memory session store: { session_id: deque([{"role": ..., "content": ...}]) }
session_store: dict[str, deque] = {}

def get_session_history(session_id: str) -> list[dict]:
    """Return the chat history list for a given session."""
    if session_id not in session_store:
        session_store[session_id] = deque(maxlen=MAX_HISTORY)
    return list(session_store[session_id])

def append_to_session(session_id: str, role: str, content: str):
    """Append a new message to the session history."""
    if session_id not in session_store:
        session_store[session_id] = deque(maxlen=MAX_HISTORY)
    session_store[session_id].append({"role": role, "content": content})

# ─────────────────────────────────────────────
# UTILITY: Embeddings & Vector Store
# ─────────────────────────────────────────────

# Initialize embedding model once (loaded from HuggingFace, cached locally)
logger.info(f"Loading embedding model: {EMBEDDING_MODEL}")
embeddings = HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL)

# Persistent ChromaDB vector store
vector_store = Chroma(
    persist_directory=str(CHROMA_DIR),
    embedding_function=embeddings,
    collection_name="mini_notebooklm",
)

def get_retriever(k: int = TOP_K_CHUNKS):
    """Return a similarity-search retriever from the vector store."""
    return vector_store.as_retriever(search_kwargs={"k": k})

# ─────────────────────────────────────────────
# UTILITY: LLM Caller (Pluggable: Gemini / OpenAI)
# ─────────────────────────────────────────────

def call_llm(prompt: str) -> str:
    """
    Send a prompt to the configured LLM and return the response text.
    Supports: Google Gemini, OpenAI GPT.
    """
    if LLM_PROVIDER == "gemini":
        if not GEMINI_API_KEY:
            raise HTTPException(status_code=500, detail="GEMINI_API_KEY not set in .env")
        genai.configure(api_key=GEMINI_API_KEY)
        model = genai.GenerativeModel("gemini-1.5-flash")
        response = model.generate_content(prompt)
        return response.text.strip()

    elif LLM_PROVIDER == "openai":
        if not OPENAI_API_KEY:
            raise HTTPException(status_code=500, detail="OPENAI_API_KEY not set in .env")
        from openai import OpenAI
        client = OpenAI(api_key=OPENAI_API_KEY)
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.2,
        )
        return response.choices[0].message.content.strip()

    else:
        raise HTTPException(status_code=500, detail=f"Unknown LLM_PROVIDER: {LLM_PROVIDER}")

# ─────────────────────────────────────────────
# UTILITY: Document Text Extraction
# ─────────────────────────────────────────────

def extract_text_from_pdf(file_path: Path) -> str:
    """Extract all text from a PDF file."""
    reader = PdfReader(str(file_path))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(pages)

def extract_text_from_txt(file_path: Path) -> str:
    """Read raw text from a .txt file."""
    return file_path.read_text(encoding="utf-8", errors="ignore")

def extract_text_from_docx(file_path: Path) -> str:
    """Extract text paragraphs from a .docx file."""
    doc = docx.Document(str(file_path))
    return "\n".join(para.text for para in doc.paragraphs if para.text.strip())

def extract_text_from_pptx(file_path: Path) -> str:
    """Extract text from all shapes/slides in a .pptx file."""
    prs = Presentation(str(file_path))
    lines = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        lines.append(f"--- Slide {slide_num} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
    return "\n".join(lines)

def extract_text(file_path: Path) -> str:
    """Dispatch text extraction based on file extension."""
    ext = file_path.suffix.lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".txt":
        return extract_text_from_txt(file_path)
    elif ext == ".docx":
        return extract_text_from_docx(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

# ─────────────────────────────────────────────
# UTILITY: Chunking & Indexing
# ─────────────────────────────────────────────

def chunk_and_index(text: str, file_name: str) -> int:
    """
    Split text into chunks, generate embeddings, and save to ChromaDB.
    Returns the number of chunks indexed.
    """
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
        separators=["\n\n", "\n", ".", " ", ""],
    )
    chunks = splitter.split_text(text)

    # Wrap each chunk as a LangChain Document with source metadata
    documents = [
        Document(
            page_content=chunk,
            metadata={"source": file_name, "chunk_index": i},
        )
        for i, chunk in enumerate(chunks)
    ]

    # Add to ChromaDB (auto-generates embeddings via our embedding model)
    vector_store.add_documents(documents)
    logger.info(f"Indexed {len(documents)} chunks from '{file_name}'")
    return len(documents)

# ─────────────────────────────────────────────
# UTILITY: Build RAG Prompt
# ─────────────────────────────────────────────

def build_rag_prompt(question: str, context_chunks: list[str], history: list[dict]) -> str:
    """
    Construct the final prompt sent to the LLM.
    Includes chat history, retrieved context, and the current question.
    """
    context = "\n\n---\n\n".join(context_chunks)

    history_text = ""
    if history:
        history_lines = []
        for msg in history[-6:]:  # Last 3 turns (user + assistant)
            role = "User" if msg["role"] == "user" else "Assistant"
            history_lines.append(f"{role}: {msg['content']}")
        history_text = "\n".join(history_lines)

    prompt = f"""You are a helpful AI assistant. Answer the user's question using ONLY the context provided below.
If the answer is not in the context, say: "I couldn't find relevant information in the uploaded documents."
Be concise, accurate, and cite relevant parts of the context if helpful.

━━━━━━━━━━━━━ DOCUMENT CONTEXT ━━━━━━━━━━━━━
{context}

━━━━━━━━━━━━━ RECENT CONVERSATION ━━━━━━━━━━
{history_text if history_text else "(No prior conversation)"}

━━━━━━━━━━━━━ CURRENT QUESTION ━━━━━━━━━━━━
{question}

Answer:"""
    return prompt

# ─────────────────────────────────────────────
# UPLOAD ROUTE — POST /upload
# ─────────────────────────────────────────────

@app.post("/upload", summary="Upload a document (PDF, TXT, DOCX, PPTX)")
async def upload_file(file: UploadFile = File(...)):
    """
    Accepts a file upload, saves it, extracts text,
    chunks it, and indexes it into ChromaDB.
    """
    # Validate extension
    file_ext = Path(file.filename).suffix.lower()
    if file_ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file type '{file_ext}'. Allowed: {', '.join(ALLOWED_EXTENSIONS)}",
        )

    # Save file to uploads/
    save_path = UPLOAD_DIR / file.filename
    try:
        with open(save_path, "wb") as f:
            shutil.copyfileobj(file.file, f)
        logger.info(f"Saved uploaded file: {save_path}")
    except Exception as e:
        logger.error(f"Failed to save file: {e}")
        raise HTTPException(status_code=500, detail=f"File save failed: {str(e)}")

    # Extract text
    try:
        text = extract_text(save_path)
        if not text.strip():
            raise ValueError("Extracted text is empty. The file may be blank or image-only.")
    except Exception as e:
        logger.error(f"Text extraction failed: {e}")
        save_path.unlink(missing_ok=True)  # Cleanup on failure
        raise HTTPException(status_code=422, detail=f"Text extraction failed: {str(e)}")

    # Chunk & index into ChromaDB
    try:
        num_chunks = chunk_and_index(text, file.filename)
    except Exception as e:
        logger.error(f"Indexing failed: {e}")
        raise HTTPException(status_code=500, detail=f"Indexing failed: {str(e)}")

    return JSONResponse(
        status_code=200,
        content={
            "message": f"File '{file.filename}' uploaded and indexed successfully.",
            "file_name": file.filename,
            "chunks_indexed": num_chunks,
            "file_size_bytes": save_path.stat().st_size,
        },
    )

# ─────────────────────────────────────────────
# ASK ROUTE — POST /ask
# ─────────────────────────────────────────────

@app.post("/ask", response_model=AskResponse, summary="Ask a question about your documents")
async def ask_question(payload: AskRequest):
    """
    Performs semantic search over indexed documents,
    builds a context-rich prompt, and queries the LLM.
    """
    question = payload.question.strip()
    if not question:
        raise HTTPException(status_code=400, detail="Question cannot be empty.")

    session_id = payload.session_id or "default"

    # Retrieve relevant chunks from ChromaDB
    try:
        retriever = get_retriever(k=TOP_K_CHUNKS)
        relevant_docs = retriever.invoke(question)
    except Exception as e:
        logger.error(f"Retrieval error: {e}")
        raise HTTPException(status_code=500, detail=f"Document retrieval failed: {str(e)}")

    if not relevant_docs:
        return AskResponse(
            answer="No documents have been uploaded yet. Please upload a file first.",
            sources=[],
            session_id=session_id,
        )

    # Extract chunk texts and source metadata
    context_chunks = [doc.page_content for doc in relevant_docs]
    sources = [
        {
            "source": doc.metadata.get("source", "unknown"),
            "chunk_index": doc.metadata.get("chunk_index", -1),
            "preview": doc.page_content[:200] + "..." if len(doc.page_content) > 200 else doc.page_content,
        }
        for doc in relevant_docs
    ]

    # Get session history
    history = get_session_history(session_id)

    # Build and send RAG prompt
    prompt = build_rag_prompt(question, context_chunks, history)
    try:
        answer = call_llm(prompt)
    except Exception as e:
        logger.error(f"LLM call failed: {e}")
        raise HTTPException(status_code=500, detail=f"LLM error: {str(e)}")

    # Update session memory
    append_to_session(session_id, "user", question)
    append_to_session(session_id, "assistant", answer)

    logger.info(f"[{session_id}] Q: {question[:60]}... → A: {answer[:60]}...")

    return AskResponse(
        answer=answer,
        sources=sources if payload.include_sources else [],
        session_id=session_id,
    )

# ─────────────────────────────────────────────
# SUMMARIZE ROUTE — POST /summarize
# ─────────────────────────────────────────────

@app.post("/summarize", response_model=SummarizeResponse, summary="Summarize uploaded documents")
async def summarize(payload: SummarizeRequest):
    """
    Retrieves a broad set of chunks and asks the LLM for a concise summary.
    If file_name is provided, summarizes that specific file.
    """
    # Determine what to summarize
    if payload.file_name:
        # Summarize a specific file
        file_path = UPLOAD_DIR / payload.file_name
        if not file_path.exists():
            raise HTTPException(status_code=404, detail=f"File '{payload.file_name}' not found.")
        try:
            text = extract_text(file_path)
        except Exception as e:
            raise HTTPException(status_code=422, detail=f"Text extraction failed: {str(e)}")

        # Truncate to avoid token limits (~12,000 chars ≈ ~3,000 tokens)
        text_sample = text[:12000]
        label = payload.file_name
    else:
        # Summarize all indexed content
        all_docs = vector_store.get()
        if not all_docs or not all_docs.get("documents"):
            raise HTTPException(
                status_code=404,
                detail="No documents indexed yet. Please upload a file first.",
            )
        # Join a sample of all chunks
        text_sample = "\n\n".join(all_docs["documents"][:30])[:12000]
        label = "all uploaded documents"

    prompt = f"""You are a document summarization expert.
Provide a clear, structured, and concise summary of the following document content.
Include: main topics, key points, and any notable conclusions.
Keep the summary under 300 words.

DOCUMENT CONTENT:
{text_sample}

SUMMARY:"""

    try:
        summary = call_llm(prompt)
    except Exception as e:
        logger.error(f"Summarization LLM call failed: {e}")
        raise HTTPException(status_code=500, detail=f"LLM error: {str(e)}")

    logger.info(f"Summarized: {label}")
    return SummarizeResponse(summary=summary, file_name=payload.file_name)

# ─────────────────────────────────────────────
# LIST FILES ROUTE — GET /files
# ─────────────────────────────────────────────

@app.get("/files", summary="List all uploaded files")
async def list_files():
    """
    Returns a list of all files currently in the uploads directory
    along with their sizes and upload timestamps.
    """
    files = []
    for f in sorted(UPLOAD_DIR.iterdir()):
        if f.is_file():
            stat = f.stat()
            files.append({
                "file_name": f.name,
                "size_bytes": stat.st_size,
                "uploaded_at": datetime.fromtimestamp(stat.st_mtime).isoformat(),
                "extension": f.suffix.lower(),
            })
    return {"files": files, "total": len(files)}

# ─────────────────────────────────────────────
# DELETE FILE ROUTE — DELETE /files/{file_name}
# ─────────────────────────────────────────────

@app.delete("/files/{file_name}", response_model=DeleteResponse, summary="Delete an uploaded file")
async def delete_file(file_name: str):
    """
    Deletes a file from the uploads directory.
    Note: ChromaDB chunks from this file remain indexed.
    To purge from the vector store, restart and re-index.
    """
    file_path = UPLOAD_DIR / file_name
    if not file_path.exists():
        raise HTTPException(status_code=404, detail=f"File '{file_name}' not found.")

    try:
        file_path.unlink()
        logger.info(f"Deleted file: {file_name}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to delete file: {str(e)}")

    return DeleteResponse(
        message=f"File '{file_name}' deleted successfully.",
        file_name=file_name,
    )

# ─────────────────────────────────────────────
# SESSION HISTORY ROUTE — GET /session/{session_id}
# ─────────────────────────────────────────────

@app.get("/session/{session_id}", summary="Get chat history for a session")
async def get_session(session_id: str):
    """Returns the current chat history for a given session ID."""
    history = get_session_history(session_id)
    return {
        "session_id": session_id,
        "messages": history,
        "count": len(history),
    }

@app.delete("/session/{session_id}", summary="Clear chat history for a session")
async def clear_session(session_id: str):
    """Clears the chat history for a given session ID."""
    if session_id in session_store:
        del session_store[session_id]
    return {"message": f"Session '{session_id}' cleared.", "session_id": session_id}

# ─────────────────────────────────────────────
# HEALTH CHECK ROUTE — GET /health
# ─────────────────────────────────────────────

@app.get("/health", response_model=HealthResponse, summary="Server health check")
async def health_check():
    """Returns the current server status, LLM provider, and upload count."""
    uploaded_count = sum(1 for f in UPLOAD_DIR.iterdir() if f.is_file())
    return HealthResponse(
        status="ok",
        provider=LLM_PROVIDER,
        uploaded_files=uploaded_count,
        timestamp=datetime.utcnow().isoformat() + "Z",
    )

# Root redirect to docs
@app.get("/", include_in_schema=False)
async def root():
    return JSONResponse(
        content={
            "message": "Mini NotebookLM API is running!",
            "docs": "http://localhost:5000/docs",
            "health": "http://localhost:5000/health",
        }
    )

# ─────────────────────────────────────────────
# VOICE OVERVIEW ROUTE — POST /voice-overview
# ─────────────────────────────────────────────

@app.post("/voice-overview", summary="Generate a spoken audio overview of your documents")
async def voice_overview(payload: VoiceOverviewRequest):
    """
    Summarises the requested document(s) with the LLM, converts the summary
    to speech via gTTS, saves an MP3 to voice_overviews/, and returns a
    download URL at GET /voice-overview/<filename>.
    """
    # ── Build text sample ──────────────────────────────────────────────────
    if payload.file_name:
        file_path = UPLOAD_DIR / payload.file_name
        if not file_path.exists():
            raise HTTPException(status_code=404, detail=f"File '{payload.file_name}' not found.")
        try:
            text = extract_text(file_path)
        except Exception as e:
            raise HTTPException(status_code=422, detail=f"Text extraction failed: {str(e)}")
        text_sample = text[:12000]
        label = payload.file_name
    else:
        all_docs = vector_store.get()
        if not all_docs or not all_docs.get("documents"):
            raise HTTPException(
                status_code=404,
                detail="No documents indexed yet. Please upload a file first.",
            )
        text_sample = "\n\n".join(all_docs["documents"][:30])[:12000]
        label = "all uploaded documents"

    # ── Ask LLM for a spoken-style summary ────────────────────────────────
    prompt = f"""You are creating a friendly spoken-audio overview for a user.
Write a clear, engaging, conversational summary of the following document content.
The summary should sound natural when read aloud — no bullet points, no markdown.
Keep it under 250 words.

DOCUMENT CONTENT:
{text_sample}

SPOKEN SUMMARY:"""

    try:
        summary_text = call_llm(prompt)
    except Exception as e:
        logger.error(f"Voice-overview LLM call failed: {e}")
        raise HTTPException(status_code=500, detail=f"LLM error: {str(e)}")

    # ── Convert to speech ─────────────────────────────────────────────────
    audio_filename = f"overview_{uuid.uuid4().hex[:8]}.mp3"
    audio_path = VOICE_DIR / audio_filename
    try:
        tts = gTTS(text=summary_text, lang=payload.language or "en", slow=False)
        tts.save(str(audio_path))
        logger.info(f"Saved voice overview: {audio_path}")
    except Exception as e:
        logger.error(f"gTTS failed: {e}")
        raise HTTPException(status_code=500, detail=f"Text-to-speech conversion failed: {str(e)}")

    return {
        "message": "Voice overview generated successfully.",
        "audio_file": audio_filename,
        "download_url": f"/voice-overview/{audio_filename}",
        "summary_text": summary_text,
        "file_name": payload.file_name,
    }


@app.get("/voice-overview/{filename}", summary="Download a generated voice overview MP3")
async def download_voice_overview(filename: str):
    """Serves the generated MP3 file for download / streaming."""
    audio_path = VOICE_DIR / filename
    if not audio_path.exists():
        raise HTTPException(status_code=404, detail=f"Audio file '{filename}' not found.")
    return FileResponse(
        path=str(audio_path),
        media_type="audio/mpeg",
        filename=filename,
    )


# ─────────────────────────────────────────────
# MIND MAP ROUTE — POST /mindmap
# ─────────────────────────────────────────────

@app.post("/mindmap", response_model=MindMapResponse, summary="Generate a mind map from documents")
async def generate_mindmap(payload: MindMapRequest):
    """
    Analyses the requested document(s) with the LLM and produces either:
    - A Mermaid.js mindmap diagram string  (output_format="mermaid")
    - A structured JSON hierarchy          (output_format="json")
    """
    # ── Build text sample ──────────────────────────────────────────────────
    if payload.file_name:
        file_path = UPLOAD_DIR / payload.file_name
        if not file_path.exists():
            raise HTTPException(status_code=404, detail=f"File '{payload.file_name}' not found.")
        try:
            text = extract_text(file_path)
        except Exception as e:
            raise HTTPException(status_code=422, detail=f"Text extraction failed: {str(e)}")
        text_sample = text[:10000]
        label = payload.file_name
    else:
        all_docs = vector_store.get()
        if not all_docs or not all_docs.get("documents"):
            raise HTTPException(
                status_code=404,
                detail="No documents indexed yet. Please upload a file first.",
            )
        text_sample = "\n\n".join(all_docs["documents"][:25])[:10000]
        label = "all uploaded documents"

    output_format = (payload.output_format or "mermaid").lower()

    # ── Build format-specific prompt ──────────────────────────────────────
    if output_format == "json":
        prompt = f"""You are a knowledge-graph expert.
Analyse the document content below and produce a hierarchical mind-map as valid JSON.
The JSON must follow this exact schema:
{{
  "root": "<central topic>",
  "branches": [
    {{
      "topic": "<branch name>",
      "subtopics": ["<item1>", "<item2>", ...]
    }},
    ...
  ]
}}
Return ONLY the raw JSON — no markdown fences, no extra text.

DOCUMENT CONTENT:
{text_sample}

JSON MIND MAP:"""
    else:  # mermaid (default)
        prompt = f"""You are a knowledge-graph expert.
Analyse the document content below and produce a Mermaid.js mindmap diagram.
Use the mindmap syntax:

mindmap
  root((Central Topic))
    Branch A
      Sub-item 1
      Sub-item 2
    Branch B
      Sub-item 3

Rules:
- Indent with 2 spaces per level.
- Keep node labels concise (under 6 words).
- Include 3–6 top-level branches and 2–4 sub-items each.
- Return ONLY the raw Mermaid diagram text — no markdown fences, no extra text.

DOCUMENT CONTENT:
{text_sample}

MERMAID MIND MAP:"""

    try:
        mind_map = call_llm(prompt)
    except Exception as e:
        logger.error(f"Mindmap LLM call failed: {e}")
        raise HTTPException(status_code=500, detail=f"LLM error: {str(e)}")

    # Strip accidental markdown code fences if the LLM added them
    mind_map = mind_map.strip()
    for fence in ("```mermaid", "```json", "```"):
        if mind_map.startswith(fence):
            mind_map = mind_map[len(fence):].strip()
        if mind_map.endswith("```"):
            mind_map = mind_map[:-3].strip()

    logger.info(f"Mind map generated for: {label} (format={output_format})")
    return MindMapResponse(
        mind_map=mind_map,
        output_format=output_format,
        file_name=payload.file_name,
    )


# ─────────────────────────────────────────────
# RUN SERVER
# ─────────────────────────────────────────────

if __name__ == "__main__":
    logger.info("━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━")
    logger.info("  Mini NotebookLM — Backend Starting Up  ")
    logger.info(f"  LLM Provider : {LLM_PROVIDER.upper()}")
    logger.info(f"  Embedding    : {EMBEDDING_MODEL}")
    logger.info(f"  Upload Dir   : {UPLOAD_DIR.resolve()}")
    logger.info(f"  Vector DB    : {CHROMA_DIR.resolve()}")
    logger.info(f"  Server URL   : http://{HOST}:{PORT}")
    logger.info(f"  API Docs     : http://localhost:{PORT}/docs")
    logger.info("━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━")
    uvicorn.run(
        "main:app",
        host=HOST,
        port=PORT,
        reload=True,     # Hot-reload on file changes (dev mode)
        log_level="info",
    )
